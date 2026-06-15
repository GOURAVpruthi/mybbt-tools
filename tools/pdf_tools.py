"""
PDF Tools — Complete ilovepdf-style toolkit.
Tools: Compress, Merge, Split, Rotate, PDF→Images, Images→PDF,
       Page Numbers, Protect, Unlock, Extract Pages, Organize Pages.
Uses: PyMuPDF (fitz), pikepdf, Pillow — all already installed.
"""

import os
import io
import shutil
import zipfile
from datetime import datetime
from pathlib import Path


class PDFTools:
    def __init__(self, upload_folder, output_folder):
        self.upload_folder = upload_folder
        self.output_folder = output_folder

    # ══════════════════════════════════════════════════════════
    #  HELPER
    # ══════════════════════════════════════════════════════════
    def _ts(self):
        return datetime.now().strftime('%Y%m%d_%H%M%S')

    def _fmt(self, size_bytes):
        if size_bytes < 1024:       return f'{size_bytes} B'
        if size_bytes < 1024**2:    return f'{size_bytes/1024:.1f} KB'
        return f'{size_bytes/1024**2:.2f} MB'

    def _out(self, name):
        return os.path.join(self.output_folder, name)

    def _page_count(self, path):
        try:
            import fitz
            d = fitz.open(path)
            n = len(d); d.close(); return n
        except Exception:
            return 0

    # ══════════════════════════════════════════════════════════
    #  1. COMPRESS
    # ══════════════════════════════════════════════════════════
    def compress(self, input_path, quality='medium'):
        ts = self._ts()
        out_name = f'compressed_{ts}.pdf'
        out_path = self._out(out_name)
        try:
            original_size = os.path.getsize(input_path)
            target_mb = {
                'high':   original_size / 1024 / 1024 * 0.85,
                'medium': original_size / 1024 / 1024 * 0.50,
                'low':    original_size / 1024 / 1024 * 0.20,
            }.get(quality, original_size / 1024 / 1024 * 0.50)
            allow_text_loss = (quality == 'low')
            logs = []
            final_size, strategy, target_hit = self._smart_compress(
                input_path, out_path, target_mb, allow_text_loss, lambda m, l='info': logs.append(m))
            reduction = ((original_size - final_size) / original_size * 100) if original_size > 0 else 0
            return {
                'success': True, 'filename': out_name,
                'original_size': original_size, 'compressed_size': final_size,
                'original_size_str': self._fmt(original_size),
                'compressed_size_str': self._fmt(final_size),
                'reduction': round(max(reduction, 0), 1), 'strategy': strategy,
                'message': f'Compressed! Size reduced by {max(reduction,0):.1f}%'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  2. MERGE
    # ══════════════════════════════════════════════════════════
    def merge(self, input_paths):
        try:
            import fitz
            ts = self._ts()
            out_name = f'merged_{ts}.pdf'
            out_path = self._out(out_name)
            merged = fitz.open()
            total_pages = 0
            for path in input_paths:
                doc = fitz.open(path)
                merged.insert_pdf(doc)
                total_pages += len(doc)
                doc.close()
            merged.save(out_path, garbage=3, deflate=True)
            merged.close()
            return {
                'success': True, 'filename': out_name,
                'total_pages': total_pages, 'files_merged': len(input_paths),
                'size_str': self._fmt(os.path.getsize(out_path)),
                'message': f'Merged {len(input_paths)} PDFs → {total_pages} pages'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  3. SPLIT
    # ══════════════════════════════════════════════════════════
    def split(self, input_path, mode='every', every_n=1, ranges_str='', prefix=''):
        """
        mode='every'  → split every N pages
        mode='range'  → split by page ranges e.g. "1-3,4-7,8"
        mode='single' → each page as separate PDF
        """
        try:
            import fitz
            doc = fitz.open(input_path)
            total = len(doc)
            ts = self._ts()
            base_name = prefix or 'split'
            output_files = []

            if mode == 'single' or (mode == 'every' and every_n == 1):
                for i in range(total):
                    new_doc = fitz.open()
                    new_doc.insert_pdf(doc, from_page=i, to_page=i)
                    fname = f'{base_name}_page{i+1}_{ts}.pdf'
                    fpath = self._out(fname)
                    new_doc.save(fpath, garbage=3, deflate=True)
                    new_doc.close()
                    output_files.append(fname)

            elif mode == 'every':
                n = max(1, int(every_n))
                chunk = 1
                for start in range(0, total, n):
                    end = min(start + n - 1, total - 1)
                    new_doc = fitz.open()
                    new_doc.insert_pdf(doc, from_page=start, to_page=end)
                    fname = f'{base_name}_part{chunk}_{ts}.pdf'
                    fpath = self._out(fname)
                    new_doc.save(fpath, garbage=3, deflate=True)
                    new_doc.close()
                    output_files.append(fname)
                    chunk += 1

            elif mode == 'range':
                ranges = self._parse_ranges(ranges_str, total)
                if not ranges:
                    doc.close()
                    return {'success': False, 'error': 'Invalid page ranges. Use format: 1-3, 4-7, 8'}
                for idx, (start, end) in enumerate(ranges, 1):
                    new_doc = fitz.open()
                    new_doc.insert_pdf(doc, from_page=start, to_page=end)
                    pages_label = f'p{start+1}-{end+1}' if start != end else f'p{start+1}'
                    fname = f'{base_name}_{pages_label}_{ts}.pdf'
                    fpath = self._out(fname)
                    new_doc.save(fpath, garbage=3, deflate=True)
                    new_doc.close()
                    output_files.append(fname)

            doc.close()

            # If multiple files → zip them
            if len(output_files) > 1:
                zip_name = f'split_result_{ts}.zip'
                zip_path = self._out(zip_name)
                with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
                    for fname in output_files:
                        zf.write(self._out(fname), fname)
                # Cleanup individual files
                for fname in output_files:
                    try: os.remove(self._out(fname))
                    except: pass
                return {
                    'success': True, 'filename': zip_name,
                    'is_zip': True, 'parts': len(output_files),
                    'total_pages': total,
                    'message': f'Split into {len(output_files)} parts → ZIP ready',
                    'download_url': f'/api/pdf/download/{zip_name}'
                }
            elif len(output_files) == 1:
                return {
                    'success': True, 'filename': output_files[0],
                    'is_zip': False, 'parts': 1,
                    'message': 'Split complete',
                    'download_url': f'/api/pdf/download/{output_files[0]}'
                }
            else:
                return {'success': False, 'error': 'No output generated'}

        except Exception as e:
            return {'success': False, 'error': str(e)}

    def _parse_ranges(self, ranges_str, total):
        """Parse '1-3, 5, 7-10' → [(0,2), (4,4), (6,9)] (0-indexed)"""
        result = []
        for part in ranges_str.replace(' ', '').split(','):
            part = part.strip()
            if not part: continue
            try:
                if '-' in part:
                    a, b = part.split('-', 1)
                    a, b = int(a) - 1, int(b) - 1
                    a = max(0, min(a, total - 1))
                    b = max(0, min(b, total - 1))
                    if a <= b: result.append((a, b))
                else:
                    p = int(part) - 1
                    if 0 <= p < total: result.append((p, p))
            except ValueError:
                continue
        return result

    # ══════════════════════════════════════════════════════════
    #  4. ROTATE
    # ══════════════════════════════════════════════════════════
    def rotate(self, input_path, angle=90, pages='all'):
        """
        angle: 90, 180, 270
        pages: 'all' or '1,3,5' or '2-5'
        """
        try:
            import fitz
            doc = fitz.open(input_path)
            total = len(doc)
            ts = self._ts()
            out_name = f'rotated_{ts}.pdf'
            out_path = self._out(out_name)

            if pages == 'all':
                page_list = list(range(total))
            else:
                page_list = []
                for r in self._parse_ranges(pages, total):
                    page_list.extend(range(r[0], r[1] + 1))

            rotated_count = 0
            for i, page in enumerate(doc):
                if i in page_list:
                    page.set_rotation(page.rotation + angle)
                    rotated_count += 1

            doc.save(out_path, garbage=3, deflate=True)
            doc.close()
            return {
                'success': True, 'filename': out_name,
                'pages_rotated': rotated_count, 'total_pages': total,
                'message': f'Rotated {rotated_count} page(s) by {angle}°'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  5. PDF → IMAGES (JPG/PNG)
    # ══════════════════════════════════════════════════════════
    def pdf_to_images(self, input_path, fmt='jpg', dpi=150, pages='all'):
        """Convert PDF pages to images. Returns ZIP of images."""
        try:
            import fitz
            doc = fitz.open(input_path)
            total = len(doc)
            ts = self._ts()

            if pages == 'all':
                page_list = list(range(total))
            else:
                page_list = []
                for r in self._parse_ranges(pages, total):
                    page_list.extend(range(r[0], r[1] + 1))

            fmt_lower = fmt.lower()
            if fmt_lower not in ('jpg', 'jpeg', 'png'):
                fmt_lower = 'jpg'
            fitz_fmt = 'jpeg' if fmt_lower in ('jpg', 'jpeg') else 'png'

            zip_name = f'pdf_images_{ts}.zip'
            zip_path = self._out(zip_name)

            with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
                for i in page_list:
                    page = doc[i]
                    mat = fitz.Matrix(dpi / 72, dpi / 72)
                    pix = page.get_pixmap(matrix=mat, alpha=False)
                    ext = 'jpg' if fitz_fmt == 'jpeg' else 'png'
                    img_bytes = pix.tobytes(fitz_fmt)
                    zf.writestr(f'page_{i+1:04d}.{ext}', img_bytes)

            doc.close()
            return {
                'success': True, 'filename': zip_name, 'is_zip': True,
                'pages_converted': len(page_list), 'total_pages': total,
                'format': fmt_lower.upper(), 'dpi': dpi,
                'message': f'Converted {len(page_list)} page(s) to {fmt_lower.upper()} images'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  6. IMAGES → PDF
    # ══════════════════════════════════════════════════════════
    def images_to_pdf(self, image_paths, page_size='auto'):
        """Convert images (JPG/PNG/BMP etc.) to a single PDF."""
        try:
            import fitz
            ts = self._ts()
            out_name = f'images_to_pdf_{ts}.pdf'
            out_path = self._out(out_name)
            doc = fitz.open()

            for img_path in image_paths:
                img = fitz.open(img_path)
                # Convert image to PDF page
                rect = img[0].rect if img.page_count > 0 else fitz.Rect(0, 0, 595, 842)
                img.close()

                # Use fitz to insert image as page
                img_doc = fitz.open(img_path)
                if img_doc.is_pdf:
                    doc.insert_pdf(img_doc)
                else:
                    # It's an image — convert to PDF bytes first
                    img_pdf_bytes = img_doc.convert_to_pdf()
                    img_doc.close()
                    tmp = fitz.open('pdf', img_pdf_bytes)
                    doc.insert_pdf(tmp)
                    tmp.close()
                    continue
                img_doc.close()

            doc.save(out_path, garbage=3, deflate=True)
            total_pages = len(doc)
            doc.close()

            return {
                'success': True, 'filename': out_name,
                'images_converted': len(image_paths), 'total_pages': total_pages,
                'size_str': self._fmt(os.path.getsize(out_path)),
                'message': f'Converted {len(image_paths)} image(s) to PDF ({total_pages} pages)'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  7. ADD PAGE NUMBERS
    # ══════════════════════════════════════════════════════════
    def add_page_numbers(self, input_path, position='bottom-center',
                          start_from=1, font_size=10, prefix='', suffix=''):
        """Add page numbers to PDF."""
        try:
            import fitz
            doc = fitz.open(input_path)
            total = len(doc)
            ts = self._ts()
            out_name = f'numbered_{ts}.pdf'
            out_path = self._out(out_name)

            for i, page in enumerate(doc):
                page_num = i + start_from
                text = f'{prefix}{page_num}{suffix}'
                rect = page.rect
                margin = 20

                # Position mapping
                pos_map = {
                    'bottom-center': fitz.Point(rect.width / 2, rect.height - margin),
                    'bottom-left':   fitz.Point(margin + 10, rect.height - margin),
                    'bottom-right':  fitz.Point(rect.width - margin - 20, rect.height - margin),
                    'top-center':    fitz.Point(rect.width / 2, margin + font_size),
                    'top-left':      fitz.Point(margin + 10, margin + font_size),
                    'top-right':     fitz.Point(rect.width - margin - 20, margin + font_size),
                }
                pos = pos_map.get(position, pos_map['bottom-center'])

                page.insert_text(
                    pos, text,
                    fontsize=font_size,
                    color=(0.3, 0.3, 0.3),
                    fontname='helv',
                )

            doc.save(out_path, garbage=3, deflate=True)
            doc.close()
            return {
                'success': True, 'filename': out_name,
                'total_pages': total,
                'message': f'Page numbers added to {total} pages (starting from {start_from})'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  8. PROTECT PDF (Add Password)
    # ══════════════════════════════════════════════════════════
    def protect(self, input_path, password, owner_password=None,
                allow_print=True, allow_copy=False):
        """Add password protection to PDF."""
        try:
            import pikepdf
            ts = self._ts()
            out_name = f'protected_{ts}.pdf'
            out_path = self._out(out_name)

            permissions = pikepdf.Permissions(
                extract=allow_copy,
                modify_annotation=False,
                modify_assembly=False,
                modify_form=False,
                modify_other=False,
                print_lowres=allow_print,
                print_highres=allow_print,
            )
            owner_pass = owner_password or (password + '_owner')

            pdf = pikepdf.Pdf.open(input_path)
            pdf.save(
                out_path,
                encryption=pikepdf.Encryption(
                    user=password,
                    owner=owner_pass,
                    allow=permissions,
                )
            )
            pdf.close()
            return {
                'success': True, 'filename': out_name,
                'message': f'PDF protected with password successfully!'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  9. UNLOCK PDF — 4-method approach (ilovepdf level)
    # ══════════════════════════════════════════════════════════
    def unlock(self, input_path, password=''):
        """Remove ALL restrictions from PDF using 4 progressive methods."""
        ts = self._ts()
        out_name = f'unlocked_{ts}.pdf'
        out_path = self._out(out_name)

        # ── Method 1: pikepdf — open with password, save without encryption
        try:
            import pikepdf
            pdf = pikepdf.Pdf.open(input_path, password=password,
                                   suppress_warnings=True)
            # Remove ALL encryption + restrictions
            pdf.save(out_path, encryption=False)
            pdf.close()
            if os.path.exists(out_path) and os.path.getsize(out_path) > 0:
                return {'success': True, 'filename': out_name,
                        'method': 'pikepdf', 'message': 'PDF fully unlocked! All restrictions removed.'}
        except Exception:
            pass

        # ── Method 2: pikepdf — try without any password (owner-restriction only)
        if password:
            try:
                import pikepdf
                pdf = pikepdf.Pdf.open(input_path, suppress_warnings=True)
                pdf.save(out_path, encryption=False)
                pdf.close()
                if os.path.exists(out_path) and os.path.getsize(out_path) > 0:
                    return {'success': True, 'filename': out_name,
                            'method': 'pikepdf-nopass', 'message': 'PDF unlocked! Owner restrictions removed.'}
            except Exception:
                pass

        # ── Method 3: PyMuPDF — authenticate and re-save
        try:
            import fitz
            doc = fitz.open(input_path)
            if doc.is_encrypted:
                ok = doc.authenticate(password) if password else doc.authenticate('')
                if not ok:
                    doc.authenticate(password)  # try once more
            new_doc = fitz.open()
            for page in doc:
                new_doc.insert_pdf(doc, from_page=page.number, to_page=page.number)
            new_doc.save(out_path, garbage=4, deflate=True, encryption=fitz.PDF_ENCRYPT_NONE)
            new_doc.close(); doc.close()
            if os.path.exists(out_path) and os.path.getsize(out_path) > 0:
                return {'success': True, 'filename': out_name,
                        'method': 'fitz', 'message': 'PDF unlocked via PyMuPDF!'}
        except Exception:
            pass

        # ── Method 4: Re-render (nuclear option — works on ANY PDF)
        # Renders each page as high-quality image → new PDF
        # Loses text selectability but removes ALL restrictions
        try:
            import fitz
            doc = fitz.open(input_path)
            if doc.is_encrypted:
                doc.authenticate(password or '')
            new_doc = fitz.open()
            for page in doc:
                # Render at 150 DPI — good quality, reasonable size
                pix = page.get_pixmap(dpi=150, alpha=False)
                img_bytes = pix.tobytes('jpeg', jpg_quality=92)
                new_page = new_doc.new_page(width=page.rect.width, height=page.rect.height)
                new_page.insert_image(new_page.rect, stream=img_bytes)
            new_doc.save(out_path, garbage=4, deflate=True)
            new_doc.close(); doc.close()
            if os.path.exists(out_path) and os.path.getsize(out_path) > 0:
                return {'success': True, 'filename': out_name,
                        'method': 're-render',
                        'message': 'PDF unlocked via re-render. All restrictions removed. '
                                   '(Note: text may not be selectable in re-rendered output)'}
        except Exception as e:
            return {'success': False, 'error': f'Could not unlock PDF: {str(e)}. '
                                               f'If password-protected, provide the correct password.'}

    # ══════════════════════════════════════════════════════════
    #  10. EXTRACT PAGES
    # ══════════════════════════════════════════════════════════
    def extract_pages(self, input_path, pages_str):
        """Extract specific pages from PDF. pages_str: '1,3,5-8'"""
        try:
            import fitz
            doc = fitz.open(input_path)
            total = len(doc)
            ts = self._ts()
            out_name = f'extracted_{ts}.pdf'
            out_path = self._out(out_name)

            ranges = self._parse_ranges(pages_str, total)
            if not ranges:
                doc.close()
                return {'success': False, 'error': 'Invalid page specification. Use: 1,3,5-8'}

            new_doc = fitz.open()
            extracted = 0
            for start, end in ranges:
                new_doc.insert_pdf(doc, from_page=start, to_page=end)
                extracted += (end - start + 1)

            new_doc.save(out_path, garbage=3, deflate=True)
            new_doc.close()
            doc.close()
            return {
                'success': True, 'filename': out_name,
                'pages_extracted': extracted, 'from_total': total,
                'size_str': self._fmt(os.path.getsize(out_path)),
                'message': f'Extracted {extracted} page(s) from {total} page document'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  11. ORGANIZE / REORDER PAGES
    # ══════════════════════════════════════════════════════════
    def organize(self, input_path, order_str):
        """
        Reorder/delete pages. order_str: '3,1,2,4' (1-indexed).
        Pages not mentioned are removed.
        """
        try:
            import fitz
            doc = fitz.open(input_path)
            total = len(doc)
            ts = self._ts()
            out_name = f'organized_{ts}.pdf'
            out_path = self._out(out_name)

            # Parse order
            try:
                order = [int(x.strip()) - 1 for x in order_str.split(',') if x.strip()]
                order = [p for p in order if 0 <= p < total]
            except ValueError:
                doc.close()
                return {'success': False, 'error': 'Invalid page order. Use comma-separated numbers: 3,1,2,4'}

            if not order:
                doc.close()
                return {'success': False, 'error': 'No valid pages specified'}

            new_doc = fitz.open()
            for p in order:
                new_doc.insert_pdf(doc, from_page=p, to_page=p)

            new_doc.save(out_path, garbage=3, deflate=True)
            new_doc.close()
            doc.close()
            return {
                'success': True, 'filename': out_name,
                'pages_output': len(order), 'original_pages': total,
                'message': f'Pages reorganized: {len(order)} page(s) in new order'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  12. WATERMARK
    # ══════════════════════════════════════════════════════════
    def watermark(self, input_path, text, opacity=0.3, angle=45,
                  font_size=48, color='gray', pages='all'):
        """Add text watermark to PDF pages."""
        try:
            import fitz
            doc = fitz.open(input_path)
            total = len(doc)
            ts = self._ts()
            out_name = f'watermarked_{ts}.pdf'
            out_path = self._out(out_name)

            color_map = {
                'gray':  (0.5, 0.5, 0.5),
                'red':   (0.8, 0.1, 0.1),
                'blue':  (0.1, 0.2, 0.8),
                'green': (0.1, 0.6, 0.1),
            }
            rgb = color_map.get(color, (0.5, 0.5, 0.5))

            if pages == 'all':
                page_list = list(range(total))
            else:
                page_list = []
                for r in self._parse_ranges(pages, total):
                    page_list.extend(range(r[0], r[1] + 1))

            for i, page in enumerate(doc):
                if i not in page_list:
                    continue
                rect = page.rect
                center = fitz.Point(rect.width / 2, rect.height / 2)
                page.insert_text(
                    center, text,
                    fontsize=font_size,
                    color=rgb,
                    rotate=angle,
                    opacity=opacity,
                    fontname='helv',
                    render_mode=3,
                )

            doc.save(out_path, garbage=3, deflate=True)
            doc.close()
            return {
                'success': True, 'filename': out_name,
                'message': f'Watermark "{text}" added to {len(page_list)} page(s)'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  13. EXTRACT TEXT
    # ══════════════════════════════════════════════════════════
    def extract_text(self, input_path, output_format='txt', pages='all'):
        """Extract all text from PDF → TXT or DOCX."""
        try:
            import fitz
            doc = fitz.open(input_path)
            total = len(doc)
            ts = self._ts()

            if pages == 'all':
                page_list = list(range(total))
            else:
                page_list = []
                for r in self._parse_ranges(pages, total):
                    page_list.extend(range(r[0], r[1] + 1))

            all_text = []
            for i in page_list:
                page = doc[i]
                text = page.get_text('text')
                all_text.append(f'--- Page {i+1} ---\n{text}')

            doc.close()
            full_text = '\n\n'.join(all_text)

            if output_format == 'txt':
                out_name = f'extracted_text_{ts}.txt'
                out_path = self._out(out_name)
                with open(out_path, 'w', encoding='utf-8') as f:
                    f.write(full_text)
                return {
                    'success': True, 'filename': out_name,
                    'pages': len(page_list), 'chars': len(full_text),
                    'message': f'Text extracted from {len(page_list)} page(s) — {len(full_text):,} characters'
                }
            else:
                # DOCX output
                try:
                    from docx import Document
                    out_name = f'extracted_text_{ts}.docx'
                    out_path = self._out(out_name)
                    doc_out = Document()
                    doc_out.add_heading('Extracted Text', 0)
                    for block in all_text:
                        for line in block.split('\n'):
                            if line.startswith('---'):
                                doc_out.add_heading(line.strip('- '), level=2)
                            elif line.strip():
                                doc_out.add_paragraph(line)
                    doc_out.save(out_path)
                    return {
                        'success': True, 'filename': out_name,
                        'message': f'Text extracted to Word document — {len(page_list)} page(s)'
                    }
                except ImportError:
                    # Fallback to txt
                    out_name = f'extracted_text_{ts}.txt'
                    out_path = self._out(out_name)
                    with open(out_path, 'w', encoding='utf-8') as f:
                        f.write(full_text)
                    return {
                        'success': True, 'filename': out_name,
                        'message': f'Text extracted (TXT format, python-docx not installed)'
                    }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  14. REPAIR PDF
    # ══════════════════════════════════════════════════════════
    def repair(self, input_path):
        """Try to repair a corrupted or malformed PDF using pikepdf."""
        try:
            import pikepdf
            ts = self._ts()
            out_name = f'repaired_{ts}.pdf'
            out_path = self._out(out_name)

            pdf = pikepdf.Pdf.open(input_path, suppress_warnings=False,
                                   attempt_recovery=True)
            pdf.save(out_path, compress_streams=True,
                     object_stream_mode=pikepdf.ObjectStreamMode.generate)
            pdf.close()

            orig = os.path.getsize(input_path)
            fixed = os.path.getsize(out_path)
            return {
                'success': True, 'filename': out_name,
                'original_size_str': self._fmt(orig),
                'repaired_size_str': self._fmt(fixed),
                'message': 'PDF repaired successfully! File structure fixed.'
            }
        except Exception as e:
            return {'success': False, 'error': f'Could not repair: {str(e)}'}

    # ══════════════════════════════════════════════════════════
    #  15. FLATTEN PDF (remove form fields, annotations)
    # ══════════════════════════════════════════════════════════
    def flatten(self, input_path):
        """Flatten PDF form fields and annotations into static content."""
        try:
            import fitz
            doc = fitz.open(input_path)
            ts = self._ts()
            out_name = f'flattened_{ts}.pdf'
            out_path = self._out(out_name)

            # Flatten by converting each page to pixmap and rebuilding
            # This removes all form fields, annotations, signatures
            new_doc = fitz.open()
            for page in doc:
                # Render at 2x for quality then embed as image
                mat = fitz.Matrix(2, 2)
                pix = page.get_pixmap(matrix=mat, alpha=False)
                img_bytes = pix.tobytes('jpeg', jpg_quality=92)
                new_page = new_doc.new_page(width=page.rect.width, height=page.rect.height)
                new_page.insert_image(new_page.rect, stream=img_bytes)

            new_doc.save(out_path, garbage=3, deflate=True)
            new_doc.close()
            doc.close()

            return {
                'success': True, 'filename': out_name,
                'message': 'PDF flattened! All form fields and annotations removed.'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  16. PDF → WORD (.docx)
    # ══════════════════════════════════════════════════════════
    def pdf_to_word(self, input_path):
        """Convert PDF to Word document using pdf2docx."""
        try:
            from pdf2docx import Converter
            ts = self._ts()
            out_name = f'converted_{ts}.docx'
            out_path = self._out(out_name)

            cv = Converter(input_path)
            cv.convert(out_path, start=0, end=None)
            cv.close()

            return {
                'success': True, 'filename': out_name,
                'size_str': self._fmt(os.path.getsize(out_path)),
                'message': 'PDF converted to Word document (.docx)'
            }
        except ImportError:
            return {'success': False,
                    'error': 'pdf2docx not installed. Run: pip install pdf2docx'}
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  17. PDF → EXCEL (.xlsx)
    # ══════════════════════════════════════════════════════════
    def pdf_to_excel(self, input_path):
        """Extract tables from PDF → Excel using pdfplumber."""
        try:
            import pdfplumber
            import openpyxl
            ts = self._ts()
            out_name = f'converted_{ts}.xlsx'
            out_path = self._out(out_name)

            wb = openpyxl.Workbook()
            wb.remove(wb.active)
            total_tables = 0

            with pdfplumber.open(input_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    tables = page.extract_tables()
                    if not tables:
                        # No tables — extract text into sheet
                        text = page.extract_text() or ''
                        if text.strip():
                            ws = wb.create_sheet(title=f'Page{page_num}_Text')
                            for line in text.split('\n'):
                                ws.append([line])
                        continue

                    for t_idx, table in enumerate(tables, 1):
                        total_tables += 1
                        sheet_name = f'P{page_num}_T{t_idx}'[:31]
                        ws = wb.create_sheet(title=sheet_name)

                        # Header styling
                        from openpyxl.styles import Font, PatternFill, Alignment
                        header_fill = PatternFill('solid', fgColor='4472C4')
                        header_font = Font(color='FFFFFF', bold=True)

                        for row_idx, row in enumerate(table, 1):
                            for col_idx, cell_val in enumerate(row or [], 1):
                                cell = ws.cell(row=row_idx, column=col_idx,
                                               value=str(cell_val or '').strip())
                                if row_idx == 1:
                                    cell.fill = header_fill
                                    cell.font = header_font
                                    cell.alignment = Alignment(horizontal='center')
                        # Auto column width
                        for col in ws.columns:
                            max_w = max((len(str(c.value or '')) for c in col), default=10)
                            ws.column_dimensions[col[0].column_letter].width = min(max_w + 2, 50)

            if not wb.sheetnames:
                wb.create_sheet('No Tables Found')
                wb.active.append(['No tables were detected in this PDF.'])

            wb.save(out_path)
            return {
                'success': True, 'filename': out_name,
                'tables_found': total_tables,
                'size_str': self._fmt(os.path.getsize(out_path)),
                'message': f'Found {total_tables} table(s) → exported to Excel'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  18. OFFICE → PDF (Word / Excel / PPT)
    #  Uses MS Office via win32com (Windows only, needs Office installed)
    # ══════════════════════════════════════════════════════════
    def office_to_pdf(self, input_path):
        """Convert Word/Excel/PPT to PDF using MS Office COM automation."""
        ext = os.path.splitext(input_path)[1].lower()
        try:
            import win32com.client
            import pythoncom
            pythoncom.CoInitialize()
            ts = self._ts()
            out_name = f'converted_{ts}.pdf'
            out_path = self._out(out_name)
            abs_in  = os.path.abspath(input_path)
            abs_out = os.path.abspath(out_path)

            if ext in ('.doc', '.docx', '.odt', '.rtf'):
                app = win32com.client.Dispatch('Word.Application')
                app.Visible = False
                try:
                    doc = app.Documents.Open(abs_in)
                    doc.SaveAs(abs_out, FileFormat=17)  # 17 = PDF
                    doc.Close()
                finally:
                    app.Quit()

            elif ext in ('.xls', '.xlsx', '.ods', '.csv'):
                app = win32com.client.Dispatch('Excel.Application')
                app.Visible = False
                try:
                    wb = app.Workbooks.Open(abs_in)
                    wb.ExportAsFixedFormat(0, abs_out)  # 0 = PDF
                    wb.Close(False)
                finally:
                    app.Quit()

            elif ext in ('.ppt', '.pptx', '.odp'):
                app = win32com.client.Dispatch('PowerPoint.Application')
                try:
                    prs = app.Presentations.Open(abs_in, WithWindow=False)
                    prs.SaveAs(abs_out, 32)  # 32 = PDF
                    prs.Close()
                finally:
                    app.Quit()
            else:
                return {'success': False, 'error': f'Unsupported format: {ext}'}

            pythoncom.CoUninitialize()

            if not os.path.exists(out_path):
                return {'success': False, 'error': 'Conversion failed — output not created'}

            return {
                'success': True, 'filename': out_name,
                'size_str': self._fmt(os.path.getsize(out_path)),
                'source_format': ext.lstrip('.').upper(),
                'message': f'{ext.lstrip(".").upper()} converted to PDF successfully!'
            }
        except ImportError:
            # Fallback: try LibreOffice if available
            return self._office_to_pdf_libreoffice(input_path)
        except Exception as e:
            return {'success': False, 'error': f'Office conversion failed: {str(e)}. Make sure MS Office is installed.'}

    def _office_to_pdf_libreoffice(self, input_path):
        """Fallback: use LibreOffice for office→pdf if win32com not available."""
        import subprocess
        ts = self._ts()
        out_name = f'converted_{ts}.pdf'
        out_path = self._out(out_name)
        try:
            result = subprocess.run(
                ['soffice', '--headless', '--convert-to', 'pdf',
                 '--outdir', self.output_folder, input_path],
                capture_output=True, timeout=60
            )
            # LibreOffice names it based on input filename
            base = os.path.splitext(os.path.basename(input_path))[0] + '.pdf'
            lo_out = os.path.join(self.output_folder, base)
            if os.path.exists(lo_out):
                os.rename(lo_out, out_path)
                return {
                    'success': True, 'filename': out_name,
                    'message': 'Converted using LibreOffice'
                }
            return {'success': False, 'error': 'LibreOffice conversion failed. Install MS Office or LibreOffice.'}
        except Exception as e:
            return {'success': False, 'error': f'No conversion engine found. Install MS Office or LibreOffice. ({e})'}

    # ══════════════════════════════════════════════════════════
    #  19. HTML → PDF
    # ══════════════════════════════════════════════════════════
    def html_to_pdf(self, html_content=None, html_path=None, url=None):
        """Convert HTML content / file / URL to PDF."""
        ts = self._ts()
        out_name = f'html_to_pdf_{ts}.pdf'
        out_path = self._out(out_name)

        # Try weasyprint first
        try:
            from weasyprint import HTML
            if url:
                HTML(url=url).write_pdf(out_path)
            elif html_path:
                HTML(filename=html_path).write_pdf(out_path)
            elif html_content:
                HTML(string=html_content).write_pdf(out_path)
            else:
                return {'success': False, 'error': 'Provide HTML content, file, or URL'}

            return {
                'success': True, 'filename': out_name,
                'size_str': self._fmt(os.path.getsize(out_path)),
                'message': 'HTML converted to PDF successfully'
            }
        except ImportError:
            pass

        # Fallback: pdfkit
        try:
            import pdfkit
            if url:
                pdfkit.from_url(url, out_path)
            elif html_path:
                pdfkit.from_file(html_path, out_path)
            elif html_content:
                pdfkit.from_string(html_content, out_path)
            return {
                'success': True, 'filename': out_name,
                'message': 'HTML converted to PDF (via pdfkit)'
            }
        except ImportError:
            pass

        return {'success': False,
                'error': 'Install weasyprint: pip install weasyprint'}

    # ══════════════════════════════════════════════════════════
    #  PDF INFO (for UI preview)
    # ══════════════════════════════════════════════════════════
    def get_info(self, input_path):
        """Get PDF metadata and page count."""
        try:
            import fitz
            doc = fitz.open(input_path)
            meta = doc.metadata
            pages = len(doc)
            size = os.path.getsize(input_path)
            # First page dimensions
            if pages > 0:
                rect = doc[0].rect
                w, h = round(rect.width), round(rect.height)
            else:
                w = h = 0
            doc.close()
            return {
                'success': True,
                'pages': pages,
                'size_str': self._fmt(size),
                'size_bytes': size,
                'title': meta.get('title', ''),
                'author': meta.get('author', ''),
                'page_width': w,
                'page_height': h,
                'encrypted': False
            }
        except Exception as e:
            # Maybe encrypted
            try:
                import pikepdf
                pikepdf.Pdf.open(input_path)
            except pikepdf.PasswordError:
                return {'success': True, 'pages': '?', 'encrypted': True,
                        'size_str': self._fmt(os.path.getsize(input_path))}
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  COMPRESSION ENGINE (from user's original gst_tools_suite.py)
    # ══════════════════════════════════════════════════════════
    def _smart_compress(self, in_path, out_path, target_mb, allow_text_loss, log):
        if os.path.abspath(in_path) == os.path.abspath(out_path):
            raise ValueError("Input and output paths cannot be the same.")
        target_bytes = int(target_mb * 1024 * 1024)
        orig_bytes = os.path.getsize(in_path)
        if orig_bytes <= target_bytes:
            shutil.copy2(in_path, out_path)
            return orig_bytes, "already optimal", True
        best_size, best_strategy = orig_bytes, "original"

        # Tier 1: Lossless
        try:
            self._compress_lossless(in_path, out_path)
            sz = os.path.getsize(out_path)
            pct = (1 - sz / orig_bytes) * 100
            if sz <= target_bytes:
                return sz, "lossless", True
            best_size, best_strategy = sz, "lossless"
        except Exception as e:
            log(f"Lossless failed: {e}", "err")
            try: shutil.copy2(in_path, out_path)
            except: pass

        # Tier 2: Image recompress
        has_pikepdf = has_pillow = False
        try: import pikepdf; has_pikepdf = True
        except: pass
        try: from PIL import Image; has_pillow = True
        except: pass

        if has_pikepdf and has_pillow:
            for q, scale in [(80, 1.0), (70, 1.0), (60, 0.9), (50, 0.8), (40, 0.7), (30, 0.6)]:
                try:
                    tmp = out_path + ".t2.tmp.pdf"
                    self._compress_images_preserve_text(in_path, tmp, q, scale)
                    sz = os.path.getsize(tmp)
                    if sz < best_size:
                        shutil.move(tmp, out_path)
                        best_size = sz
                        best_strategy = f"image-recompress q={q}"
                    else:
                        try: os.remove(tmp)
                        except: pass
                    if sz <= target_bytes:
                        return sz, best_strategy, True
                except: continue

        if not allow_text_loss:
            return best_size, best_strategy, False

        # Tier 3: Rasterize
        try:
            import fitz
            for dpi, q in [(200, 80), (150, 75), (120, 70), (100, 60), (85, 55), (72, 50)]:
                try:
                    tmp = out_path + ".t3.tmp.pdf"
                    self._compress_rasterize(in_path, tmp, dpi, q)
                    sz = os.path.getsize(tmp)
                    if sz < best_size:
                        shutil.move(tmp, out_path)
                        best_size = sz
                        best_strategy = f"rasterized @{dpi}dpi"
                    else:
                        try: os.remove(tmp)
                        except: pass
                    if sz <= target_bytes:
                        return sz, best_strategy, True
                except: continue
        except ImportError:
            pass

        for sfx in (".t2.tmp.pdf", ".t3.tmp.pdf"):
            p = out_path + sfx
            if os.path.exists(p):
                try: os.remove(p)
                except: pass

        return best_size, best_strategy, False

    def _compress_lossless(self, in_path, out_path):
        import fitz
        doc = fitz.open(in_path)
        try:
            doc.save(out_path, garbage=4, deflate=True,
                     deflate_images=True, deflate_fonts=True, clean=True)
        finally:
            doc.close()

    def _compress_images_preserve_text(self, in_path, out_path, jpeg_quality, downscale=1.0):
        import pikepdf
        from PIL import Image
        pdf = pikepdf.Pdf.open(in_path)
        try:
            for page in pdf.pages:
                try: images = page.images
                except: continue
                for name in list(images.keys()):
                    raw_image = images[name]
                    try:
                        pdfimg = pikepdf.PdfImage(raw_image)
                        pil = pdfimg.as_pil_image()
                        if pil.mode in ("RGBA", "P", "LA"): pil = pil.convert("RGB")
                        elif pil.mode == "CMYK":            pil = pil.convert("RGB")
                        if downscale < 1.0:
                            w, h = pil.size
                            pil = pil.resize((max(1, int(w*downscale)), max(1, int(h*downscale))), Image.LANCZOS)
                        buf = io.BytesIO()
                        pil.save(buf, format="JPEG", quality=jpeg_quality, optimize=True)
                        new_bytes = buf.getvalue()
                        try: raw_size = len(bytes(raw_image.read_raw_bytes()))
                        except: raw_size = 10**9
                        if len(new_bytes) >= raw_size: continue
                        cs = pikepdf.Name("/DeviceGray") if pil.mode == "L" else pikepdf.Name("/DeviceRGB")
                        raw_image.write(new_bytes, filter=pikepdf.Name("/DCTDecode"))
                        raw_image.Width = pil.width; raw_image.Height = pil.height
                        raw_image.BitsPerComponent = 8; raw_image.ColorSpace = cs
                        for k in ("/DecodeParms", "/Decode"):
                            if k in raw_image: del raw_image[k]
                    except: continue
            pdf.save(out_path, compress_streams=True,
                     object_stream_mode=pikepdf.ObjectStreamMode.generate,
                     recompress_flate=True, linearize=False)
        finally:
            pdf.close()

    def _compress_rasterize(self, in_path, out_path, dpi, jpeg_quality):
        import fitz
        doc = fitz.open(in_path)
        new_doc = fitz.open()
        try:
            for page in doc:
                pix = page.get_pixmap(dpi=dpi, alpha=False)
                img_bytes = pix.tobytes("jpeg", jpg_quality=jpeg_quality)
                new_page = new_doc.new_page(width=page.rect.width, height=page.rect.height)
                new_page.insert_image(new_page.rect, stream=img_bytes)
            new_doc.save(out_path, garbage=4, deflate=True)
        finally:
            new_doc.close(); doc.close()

    # ══════════════════════════════════════════════════════════
    #  NEW — GET PAGE THUMBNAILS
    # ══════════════════════════════════════════════════════════
    def get_thumbnails(self, input_path, dpi=72):
        """Return base64 PNG thumbnails for every page."""
        try:
            import fitz, base64
            doc = fitz.open(input_path)
            if doc.is_encrypted:
                doc.authenticate('')
            thumbnails = []
            mat = fitz.Matrix(dpi / 72, dpi / 72)
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=mat, alpha=False)
                b64 = base64.b64encode(pix.tobytes('png')).decode()
                thumbnails.append({
                    'page': i + 1,
                    'data': f'data:image/png;base64,{b64}',
                    'width': pix.width,
                    'height': pix.height,
                    'rotation': page.rotation,
                })
            doc.close()
            return {'success': True, 'thumbnails': thumbnails, 'total': len(thumbnails)}
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  NEW — REMOVE PAGES
    # ══════════════════════════════════════════════════════════
    def remove_pages(self, input_path, pages_str):
        """Remove specific pages. pages_str: '1,3,5-8'"""
        try:
            import fitz
            doc = fitz.open(input_path)
            total = len(doc)
            ranges = self._parse_ranges(pages_str, total)
            if not ranges:
                doc.close()
                return {'success': False, 'error': 'Invalid page spec. Use: 1,3,5-8'}
            # Collect all page indices to remove (0-based), deduplicated
            to_remove = set()
            for s, e in ranges:
                to_remove.update(range(s, e + 1))
            if len(to_remove) >= total:
                doc.close()
                return {'success': False, 'error': 'Cannot remove all pages!'}
            # Delete in reverse order to preserve indices
            for idx in sorted(to_remove, reverse=True):
                doc.delete_page(idx)
            ts = self._ts()
            out_name = f'removed_{ts}.pdf'
            out_path = self._out(out_name)
            doc.save(out_path, garbage=3, deflate=True)
            remaining = len(doc)
            doc.close()
            return {
                'success': True, 'filename': out_name,
                'removed': len(to_remove), 'remaining': remaining,
                'size_str': self._fmt(os.path.getsize(out_path)),
                'message': f'Removed {len(to_remove)} page(s). {remaining} pages remaining.'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  NEW — ORGANIZE PAGES (reorder + rotate)
    # ══════════════════════════════════════════════════════════
    def organize_pages(self, input_path, new_order, rotations=None):
        """Reorder pages and apply per-page rotations.
        new_order: list of 1-based page numbers in desired order.
        rotations: dict {"1": 90, "3": 180, ...} page_num -> degrees.
        """
        try:
            import fitz
            doc = fitz.open(input_path)
            if doc.is_encrypted:
                doc.authenticate('')
            total = len(doc)
            # Validate order
            new_order = [int(p) for p in new_order if 1 <= int(p) <= total]
            if not new_order:
                doc.close()
                return {'success': False, 'error': 'Invalid page order'}
            new_doc = fitz.open()
            for page_num in new_order:
                new_doc.insert_pdf(doc, from_page=page_num - 1, to_page=page_num - 1)
                if rotations and str(page_num) in rotations:
                    rot = int(rotations[str(page_num)]) % 360
                    if rot:
                        new_doc[-1].set_rotation(rot)
            ts = self._ts()
            out_name = f'organized_{ts}.pdf'
            out_path = self._out(out_name)
            new_doc.save(out_path, garbage=3, deflate=True)
            new_doc.close(); doc.close()
            return {
                'success': True, 'filename': out_name,
                'pages': len(new_order),
                'size_str': self._fmt(os.path.getsize(out_path)),
                'message': f'PDF organized! {len(new_order)} pages saved.'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  NEW — MERGE WITH CUSTOM ORDER
    # ══════════════════════════════════════════════════════════
    def merge_ordered(self, input_paths):
        """Merge PDFs in the given order (input_paths already ordered)."""
        return self.merge(input_paths)  # Existing merge handles this

    # ══════════════════════════════════════════════════════════
    #  OCR PDF — Extract text from image-based PDFs
    # ══════════════════════════════════════════════════════════
    def ocr_pdf(self, input_path):
        """Run OCR on scanned PDF pages, return searchable text + PDF file."""
        try:
            import fitz
            doc = fitz.open(input_path)
            if doc.is_encrypted:
                doc.authenticate('')
            all_text = []
            ts = self._ts()

            for i, page in enumerate(doc):
                # Try embedded text first
                text = page.get_text().strip()
                if len(text) > 20:
                    all_text.append({'page': i + 1, 'text': text, 'method': 'embedded'})
                else:
                    # Render page as image and run OCR
                    try:
                        import pytesseract
                        from PIL import Image as PILImage
                        import io as _io
                        pix = page.get_pixmap(dpi=300, alpha=False)  # 300dpi for better OCR
                        img_bytes = pix.tobytes('png')
                        img = PILImage.open(_io.BytesIO(img_bytes))
                        ocr_text = pytesseract.image_to_string(img, lang='eng')
                        all_text.append({'page': i + 1, 'text': ocr_text.strip(), 'method': 'ocr'})
                    except Exception as ocr_err:
                        all_text.append({'page': i + 1, 'text': f'[OCR failed: {ocr_err}]', 'method': 'error'})

            doc.close()

            # Build output PDF with extracted text
            out_doc = fitz.open()
            for pt in all_text:
                page = out_doc.new_page(width=595, height=842)  # A4
                page.insert_text(
                    (50, 40),
                    f'OCR Result — Page {pt["page"]}  ({pt["method"]})',
                    fontsize=9, fontname='helv', color=(0.5, 0.5, 0.5)
                )
                page.draw_line(
                    fitz.Point(50, 52), fitz.Point(545, 52),
                    color=(0.88, 0.1, 0.17), width=0.8
                )
                y = 68
                for line in pt['text'].split('\n'):
                    while line:
                        chunk = line[:95]
                        line = line[95:]
                        if y > 810:
                            page = out_doc.new_page(width=595, height=842)
                            page.insert_text(
                                (50, 40), f'OCR Result — Page {pt["page"]} (cont.)',
                                fontsize=9, fontname='helv', color=(0.5, 0.5, 0.5)
                            )
                            y = 60
                        page.insert_text(
                            (50, y), chunk,
                            fontsize=10, fontname='helv', color=(0.05, 0.05, 0.05)
                        )
                        y += 15
                    y += 3

            out_name = f'ocr_{ts}.pdf'
            out_path = self._out(out_name)
            out_doc.save(out_path, garbage=4, deflate=True)
            out_doc.close()

            full_text = '\n\n'.join([f"--- Page {p['page']} ---\n{p['text']}" for p in all_text])
            return {
                'success': True, 'filename': out_name,
                'pages': len(all_text),
                'preview': full_text[:500],
                'char_count': len(full_text),
                'message': f'OCR complete! Extracted text from {len(all_text)} pages. Downloading as PDF.'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}


    # ══════════════════════════════════════════════════════════
    #  CROP PDF — Remove margins
    # ══════════════════════════════════════════════════════════
    def crop_pdf(self, input_path, top=0, bottom=0, left=0, right=0):
        """Crop PDF margins. Values in points (72pt = 1 inch)."""
        try:
            import fitz
            doc = fitz.open(input_path)
            if doc.is_encrypted:
                doc.authenticate('')
            for page in doc:
                r = page.rect
                new_rect = fitz.Rect(
                    r.x0 + float(left),
                    r.y0 + float(top),
                    r.x1 - float(right),
                    r.y1 - float(bottom)
                )
                if new_rect.is_valid and new_rect.width > 10 and new_rect.height > 10:
                    page.set_cropbox(new_rect)
            ts = self._ts()
            out_name = f'cropped_{ts}.pdf'
            out_path = self._out(out_name)
            doc.save(out_path, garbage=3, deflate=True)
            doc.close()
            return {
                'success': True, 'filename': out_name,
                'size_str': self._fmt(os.path.getsize(out_path)),
                'message': f'PDF cropped! Margins removed: T={top}pt L={left}pt B={bottom}pt R={right}pt'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  SIGN PDF — Embed signature image
    # ══════════════════════════════════════════════════════════
    def sign_pdf(self, input_path, signature_b64, page_num=None,
                 position='bottom-right', sig_width=200):
        """Embed a signature (base64 PNG) into PDF."""
        try:
            import fitz, base64, io as _io
            doc = fitz.open(input_path)
            if doc.is_encrypted:
                doc.authenticate('')
            total = len(doc)
            page_idx = (total - 1) if page_num is None else max(0, int(page_num) - 1)
            page = doc[page_idx]

            # Decode signature
            sig_data = base64.b64decode(signature_b64.split(',')[-1])
            sig_width = float(sig_width)
            # Get natural aspect from Pillow
            from PIL import Image as PILImage
            pimg = PILImage.open(_io.BytesIO(sig_data))
            aspect = pimg.height / pimg.width
            sig_height = sig_width * aspect

            pr = page.rect
            margin = 20
            positions = {
                'bottom-right': fitz.Rect(pr.x1 - sig_width - margin,
                                          pr.y1 - sig_height - margin,
                                          pr.x1 - margin, pr.y1 - margin),
                'bottom-left':  fitz.Rect(margin, pr.y1 - sig_height - margin,
                                          margin + sig_width, pr.y1 - margin),
                'bottom-center':fitz.Rect((pr.width - sig_width)/2,
                                          pr.y1 - sig_height - margin,
                                          (pr.width + sig_width)/2, pr.y1 - margin),
                'top-right':    fitz.Rect(pr.x1 - sig_width - margin, margin,
                                          pr.x1 - margin, margin + sig_height),
                'top-left':     fitz.Rect(margin, margin,
                                          margin + sig_width, margin + sig_height),
            }
            rect = positions.get(position, positions['bottom-right'])
            page.insert_image(rect, stream=sig_data)

            ts = self._ts()
            out_name = f'signed_{ts}.pdf'
            out_path = self._out(out_name)
            doc.save(out_path, garbage=3, deflate=True)
            doc.close()
            return {
                'success': True, 'filename': out_name,
                'size_str': self._fmt(os.path.getsize(out_path)),
                'message': f'Signature added to page {page_idx+1} ({position})!'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  REDACT PDF — Black out sensitive text
    # ══════════════════════════════════════════════════════════
    def redact_pdf(self, input_path, search_text):
        """Search for text in PDF and permanently redact (black out) it."""
        try:
            import fitz
            doc = fitz.open(input_path)
            if doc.is_encrypted:
                doc.authenticate('')
            total_found = 0
            for page in doc:
                hits = page.search_for(search_text, quads=False)
                for rect in hits:
                    # Expand slightly for better coverage
                    expanded = fitz.Rect(rect.x0 - 2, rect.y0 - 1,
                                         rect.x1 + 2, rect.y1 + 1)
                    page.add_redact_annot(expanded, fill=(0, 0, 0))
                    total_found += 1
                if hits:
                    page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)
            ts = self._ts()
            out_name = f'redacted_{ts}.pdf'
            out_path = self._out(out_name)
            doc.save(out_path, garbage=4, deflate=True)
            doc.close()
            if total_found == 0:
                return {'success': False,
                        'error': f'Text "{search_text}" not found in PDF.'}
            return {
                'success': True, 'filename': out_name,
                'found': total_found,
                'message': f'Redacted {total_found} instance(s) of "{search_text}"!'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  COMPARE PDF — Visual side-by-side diff
    # ══════════════════════════════════════════════════════════
    def compare_pdf(self, path_a, path_b):
        """Generate side-by-side comparison images for first page of each PDF."""
        try:
            import fitz, base64
            results = []
            doc_a = fitz.open(path_a)
            doc_b = fitz.open(path_b)
            pages = min(len(doc_a), len(doc_b), 10)  # Compare up to 10 pages
            for i in range(pages):
                pix_a = doc_a[i].get_pixmap(dpi=100, alpha=False)
                pix_b = doc_b[i].get_pixmap(dpi=100, alpha=False)
                b64_a = base64.b64encode(pix_a.tobytes('png')).decode()
                b64_b = base64.b64encode(pix_b.tobytes('png')).decode()
                results.append({
                    'page': i + 1,
                    'img_a': f'data:image/png;base64,{b64_a}',
                    'img_b': f'data:image/png;base64,{b64_b}',
                })
            doc_a.close()
            doc_b.close()
            return {
                'success': True,
                'pages': pages,
                'comparisons': results,
                'message': f'Comparing {pages} page(s) side by side.'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  PDF TO POWERPOINT
    # ══════════════════════════════════════════════════════════
    def to_pptx(self, input_path, dpi=150):
        """Convert PDF pages to PowerPoint (each page = one slide)."""
        try:
            import fitz
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            import io as _io

            doc = fitz.open(input_path)
            if doc.is_encrypted:
                doc.authenticate('')

            prs = Presentation()
            blank_layout = prs.slide_layouts[6]  # Blank slide

            for page in doc:
                # Get page dimensions in inches (1 pt = 1/72 inch)
                w_in = page.rect.width / 72
                h_in = page.rect.height / 72
                prs.slide_width = Inches(w_in)
                prs.slide_height = Inches(h_in)
                slide = prs.slides.add_slide(blank_layout)
                # Render page as image
                pix = page.get_pixmap(dpi=dpi, alpha=False)
                img_bytes = pix.tobytes('jpeg', jpg_quality=90)
                img_buf = _io.BytesIO(img_bytes)
                slide.shapes.add_picture(img_buf, 0, 0,
                                          width=Inches(w_in),
                                          height=Inches(h_in))

            ts = self._ts()
            out_name = f'converted_{ts}.pptx'
            out_path = self._out(out_name)
            prs.save(out_path)
            doc.close()
            return {
                'success': True, 'filename': out_name,
                'slides': len(doc),
                'size_str': self._fmt(os.path.getsize(out_path)),
                'message': f'Converted to PowerPoint! {len(doc)} slides created.'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  PDF TO PDF/A (archival)
    # ══════════════════════════════════════════════════════════
    def to_pdfa(self, input_path):
        """Convert PDF to PDF/A-1b archival format."""
        try:
            import fitz
            doc = fitz.open(input_path)
            if doc.is_encrypted:
                doc.authenticate('')
            ts = self._ts()
            out_name = f'pdfa_{ts}.pdf'
            out_path = self._out(out_name)
            # Save with PDF/A-compatible settings
            doc.save(out_path,
                     garbage=4,
                     deflate=True,
                     clean=True,
                     linear=False,
                     encryption=fitz.PDF_ENCRYPT_NONE)
            # Add PDF/A XMP metadata
            xmp = b'''<?xpacket begin='\xef\xbb\xbf' id='W5M0MpCehiHzreSzNTczkc9d'?>
<x:xmpmeta xmlns:x='adobe:ns:meta/'>
  <rdf:RDF xmlns:rdf='http://www.w3.org/1999/02/22-rdf-syntax-ns#'>
    <rdf:Description rdf:about='' xmlns:pdfaid='http://www.aiim.org/pdfa/ns/id/'>
      <pdfaid:part>1</pdfaid:part>
      <pdfaid:conformance>B</pdfaid:conformance>
    </rdf:Description>
  </rdf:RDF>
</x:xmpmeta><?xpacket end='w'?>'''
            doc2 = fitz.open(out_path)
            doc2.set_xml_metadata(xmp.decode())
            doc2.save(out_path, incremental=True, encryption=fitz.PDF_ENCRYPT_NONE)
            doc2.close()
            doc.close()
            return {
                'success': True, 'filename': out_name,
                'size_str': self._fmt(os.path.getsize(out_path)),
                'message': 'Converted to PDF/A-1b archival format!'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  AI SUMMARIZE
    # ══════════════════════════════════════════════════════════
    def ai_summarize(self, input_path, api_key, language='English'):
        """Extract text and summarize using Gemini API, output as PDF."""
        try:
            import fitz
            doc = fitz.open(input_path)
            if doc.is_encrypted:
                doc.authenticate('')
            text = ''
            for page in doc:
                text += page.get_text()
            doc.close()
            text = text[:30000]
            if len(text.strip()) < 50:
                return {'success': False, 'error': 'Could not extract text from PDF. '
                        'The PDF may be image-based (try OCR first).'}
            import google.generativeai as genai
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel('gemini-1.5-flash')
            prompt = (f"Please provide a comprehensive summary of the following document "
                      f"in {language}. Include:\n"
                      f"1. Main topic/purpose\n"
                      f"2. Key points (bullet list)\n"
                      f"3. Important conclusions\n\n"
                      f"Document text:\n{text}")
            response = model.generate_content(prompt)
            summary = response.text
            ts = self._ts()
            out_name = f'summary_{ts}.pdf'
            out_path = self._out(out_name)
            # Write summary as a clean PDF
            out_doc = fitz.open()
            page = out_doc.new_page(width=595, height=842)  # A4
            margin = 50
            page.insert_text(
                (margin, 80),
                f'AI Summary ({language})',
                fontsize=16, fontname='helv', color=(0.13, 0.13, 0.13)
            )
            page.draw_line(
                fitz.Point(margin, 100), fitz.Point(595 - margin, 100),
                color=(0.88, 0.1, 0.17), width=1.5
            )
            # Insert summary text in chunks per line
            y = 120
            for line in summary.split('\n'):
                if y > 800:
                    page = out_doc.new_page(width=595, height=842)
                    y = 50
                # Wrap long lines
                words = line if len(line) < 90 else line
                page.insert_text(
                    (margin, y), line[:110],
                    fontsize=10, fontname='helv', color=(0.1, 0.1, 0.1)
                )
                y += 15
                if len(line) > 110:
                    page.insert_text(
                        (margin, y), line[110:220],
                        fontsize=10, fontname='helv', color=(0.1, 0.1, 0.1)
                    )
                    y += 15
            out_doc.save(out_path)
            out_doc.close()
            return {
                'success': True, 'filename': out_name,
                'summary': summary,
                'word_count': len(text.split()),
                'message': f'AI summary generated successfully!'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}

    # ══════════════════════════════════════════════════════════
    #  TRANSLATE PDF
    # ══════════════════════════════════════════════════════════
    def translate_pdf(self, input_path, api_key, target_lang='Hindi'):
        """Extract text and translate using Gemini API, output as PDF."""
        try:
            import fitz
            doc = fitz.open(input_path)
            if doc.is_encrypted:
                doc.authenticate('')
            pages_text = []
            for i, page in enumerate(doc):
                t = page.get_text().strip()
                if t:
                    pages_text.append({'page': i + 1, 'text': t})
            doc.close()
            if not pages_text:
                return {'success': False,
                        'error': 'No text found in PDF. Try OCR first.'}
            import google.generativeai as genai
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel('gemini-1.5-flash')
            translated_parts = []
            for pt in pages_text:
                chunk = pt['text'][:5000]
                prompt = f"Translate the following text to {target_lang}. Return ONLY the translation, nothing else:\n\n{chunk}"
                try:
                    resp = model.generate_content(prompt)
                    translated_parts.append({'page': pt['page'], 'text': resp.text})
                except Exception:
                    translated_parts.append({'page': pt['page'], 'text': '[Translation failed for this page]'})
            ts = self._ts()
            out_name = f'translated_{ts}.pdf'
            out_path = self._out(out_name)
            # Build PDF with translated content
            out_doc = fitz.open()
            for part in translated_parts:
                page = out_doc.new_page(width=595, height=842)
                # Page header
                page.insert_text(
                    (50, 40),
                    f'Translated to {target_lang}  |  Page {part["page"]}',
                    fontsize=9, fontname='helv', color=(0.5, 0.5, 0.5)
                )
                page.draw_line(
                    fitz.Point(50, 52), fitz.Point(545, 52),
                    color=(0.88, 0.1, 0.17), width=0.8
                )
                # Insert translated text line by line
                y = 68
                for line in part['text'].split('\n'):
                    # Wrap at 95 chars
                    while line:
                        chunk = line[:95]
                        line = line[95:]
                        if y > 810:
                            page = out_doc.new_page(width=595, height=842)
                            page.insert_text(
                                (50, 40), f'Translated to {target_lang} (cont.)',
                                fontsize=9, fontname='helv', color=(0.5, 0.5, 0.5)
                            )
                            y = 60
                        page.insert_text(
                            (50, y), chunk,
                            fontsize=10, fontname='helv', color=(0.05, 0.05, 0.05)
                        )
                        y += 15
                    y += 3  # paragraph spacing
            out_doc.save(out_path)
            out_doc.close()
            full_text = '\n\n'.join([p['text'] for p in translated_parts])
            return {
                'success': True, 'filename': out_name,
                'preview': full_text[:600],
                'message': f'Translation to {target_lang} complete! Downloading as PDF.'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}
